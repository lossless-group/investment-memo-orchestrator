"""
Deck Analyst Agent - Extracts key information from pitch decks.

This agent runs FIRST if a deck is available, creating initial section drafts
that subsequent agents can build upon.
"""

from pathlib import Path
from typing import Dict, List, Any
from langchain_anthropic import ChatAnthropic
from pypdf import PdfReader
import json
import base64
import os
from anthropic import Anthropic
import fitz  # PyMuPDF
from io import BytesIO
from PIL import Image
from pptx import Presentation  # For PowerPoint files

from ..state import MemoState, DeckAnalysisData, SectionDraft
from ..outline_loader import load_outline_for_state

# Mapping from screenshot categories to section filenames in 0-deck-sections/
# Used to embed screenshots in the appropriate deck section files
SCREENSHOT_CATEGORY_TO_SECTION = {
    "team": "deck-team.md",
    "traction": "deck-traction.md",
    "product-demo": "deck-product.md",
    "technology": "deck-product.md",
    "solution": "deck-product.md",
    "market-size": "deck-market.md",
    "competitive-positioning": "deck-competitive.md",
    "competition-landscape": "deck-competitive.md",
    "unit-economics": "deck-traction.md",
    "business-model": "deck-traction.md",
    "fundraising": "deck-traction.md",
    "financials": "deck-traction.md",
    "go-to-market": "deck-traction.md",
    "branding": "deck-product.md",
    "customer-pain": "deck-product.md",
    "ideal-customer-profile": "deck-market.md",
    "overview": None,
    "problem": None,
    "value-proposition": None,
    "partnerships": None,
    "vision": None,
    "impact": None,
    "general": None,
}

# Classification guide path — loaded and passed to Claude Vision for accurate slide classification.
# See templates/deck-classification-guide.md for the full category definitions.
CLASSIFICATION_GUIDE_PATH = Path("templates/deck-classification-guide.md")

# Firm-scoped classification guide search paths (checked in order)
CLASSIFICATION_GUIDE_SEARCH_PATHS = [
    # io/{firm}/templates/deck-classification-guide.md  (resolved at runtime)
    "templates/deck-classification-guide.md",
]


def load_classification_guide(firm: str = None) -> str:
    """
    Load the deck slide classification guide markdown.

    Searches firm-scoped path first, then project-level templates/.

    Args:
        firm: Optional firm name to check io/{firm}/templates/ first

    Returns:
        Classification guide contents, or empty string if not found
    """
    search_paths = []

    if firm:
        search_paths.append(Path(f"io/{firm}/templates/deck-classification-guide.md"))

    search_paths.append(CLASSIFICATION_GUIDE_PATH)

    for path in search_paths:
        if path.exists():
            guide_text = path.read_text()
            print(f"  📋 Loaded classification guide from {path}", flush=True)
            return guide_text

    print("  ⚠️  No deck-classification-guide.md found, using inline categories", flush=True)
    return ""


# Categories considered high-signal for the Key Slides gallery in Executive Summary.
# Lower-priority categories are excluded from Key Slides to keep it focused.
KEY_SLIDES_PRIORITY_CATEGORIES = [
    "traction", "unit-economics", "team", "market-size",
    "competitive-positioning", "product-demo", "technology",
    "fundraising", "business-model",
]

# Sections that should NEVER receive deck screenshots (synthesis/summary sections)
EXCLUDED_SECTIONS = {
    "scorecard", "recommendation", "closing", "assessment",
    "investment thesis", "12ps scorecard",
}


def build_section_to_screenshots(
    screenshots: List[Dict[str, Any]],
    output_dir: Path
) -> Dict[str, Any]:
    """
    Build a mapping from screenshot categories to their primary section and metadata.

    Each screenshot is assigned to exactly ONE primary section based on its category.
    The inject_deck_images agent uses this mapping to place images, enforcing the
    rule that each image appears at most twice (once in its primary section, once
    optionally in the Executive Summary Key Slides gallery).

    Args:
        screenshots: List of screenshot metadata dicts with 'category', 'path',
                    'description', and optionally 'slug'
        output_dir: Output directory (for building paths)

    Returns:
        Dict with:
          - "images": list of image metadata dicts, each with path, category,
                      slug, description, and primary_section
          - "key_slides": list of paths for the Executive Summary Key Slides gallery
    """
    if not screenshots:
        return {"images": [], "key_slides": []}

    images = []
    key_slides_candidates = []

    for ss in screenshots:
        category = ss.get("category", "general")
        if category == "general":
            continue

        # Build path relative to project root (not absolute — stays portable)
        img_path = str(output_dir / ss["path"])
        slug = ss.get("slug", "")
        description = ss.get("description", f"Page {ss.get('page_number', '?')} from pitch deck")

        image_entry = {
            "path": img_path,
            "page_number": ss.get("page_number"),
            "category": category,
            "slug": slug,
            "description": description,
            "primary_section": category,  # inject agent maps category → section file
        }
        images.append(image_entry)

        # Collect candidates for Key Slides gallery
        if category in KEY_SLIDES_PRIORITY_CATEGORIES:
            key_slides_candidates.append(image_entry)

    # Select up to 5 key slides, prioritized by category order
    priority_order = {cat: i for i, cat in enumerate(KEY_SLIDES_PRIORITY_CATEGORIES)}
    key_slides_candidates.sort(key=lambda x: priority_order.get(x["category"], 99))
    key_slides = [entry["path"] for entry in key_slides_candidates[:5]]

    return {
        "images": images,
        "key_slides": key_slides,
    }


# Optional: pdf2image for higher quality rendering (requires Poppler)
try:
    from pdf2image import convert_from_path
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False


import concurrent.futures
import signal


def extract_screenshots_with_timeout(
    pdf_path: str,
    deck_analysis: Dict,
    output_dir: Path,
    timeout_seconds: int = 90,
    firm: str = None
) -> List[Dict[str, Any]]:
    """
    Extract visual screenshots from deck pages with a timeout.

    This function runs screenshot extraction in a separate thread with a timeout
    to prevent the workflow from hanging if Claude Vision API is slow or unresponsive.

    Args:
        pdf_path: Path to PDF file
        deck_analysis: Already extracted deck analysis data
        output_dir: Directory where screenshots should be saved
        timeout_seconds: Maximum time to wait for screenshot extraction
        firm: Optional firm name for firm-scoped classification guide

    Returns:
        List of screenshot metadata dicts, or empty list if failed/timed out
    """
    print(f"\n📷 Starting screenshot extraction (timeout: {timeout_seconds}s)...", flush=True)
    print(f"   Output directory: {output_dir}", flush=True)

    def _extract_screenshots():
        """Inner function that does the actual extraction."""
        try:
            # Create Anthropic client for vision
            client = Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))

            # Identify which pages have valuable visual content
            print("  Identifying visual pages with Claude Vision...", flush=True)
            page_selections = identify_visual_pages(pdf_path, deck_analysis, client, firm=firm)

            if not page_selections:
                print("  No significant visual content identified", flush=True)
                return []

            print(f"  Found {len(page_selections)} pages with visual content", flush=True)

            # Extract the screenshots
            screenshots = extract_deck_screenshots(
                pdf_path,
                output_dir,
                page_selections,
                use_pdf2image=PDF2IMAGE_AVAILABLE,
                dpi=150
            )

            print(f"  ✓ Extracted {len(screenshots)} screenshots", flush=True)
            return screenshots

        except Exception as e:
            print(f"  ⚠️  Screenshot extraction error: {e}", flush=True)
            return []

    # Run extraction with timeout using ThreadPoolExecutor
    try:
        with concurrent.futures.ThreadPoolExecutor(max_workers=1) as executor:
            future = executor.submit(_extract_screenshots)
            try:
                screenshots = future.result(timeout=timeout_seconds)
                return screenshots
            except concurrent.futures.TimeoutError:
                print(f"  ⚠️  Screenshot extraction timed out after {timeout_seconds}s", flush=True)
                print("  Continuing workflow without screenshots...", flush=True)
                return []
    except Exception as e:
        print(f"  ⚠️  Screenshot extraction failed: {e}", flush=True)
        print("  Continuing workflow without screenshots...", flush=True)
        return []


def embed_screenshots_in_section_files(output_dir: Path, screenshots: List[Dict[str, Any]]) -> int:
    """
    Embed screenshot images into their corresponding section files in 0-deck-sections/.

    This function updates the already-saved section files to include image embeds
    at the top of each file. Uses absolute paths for cross-directory compatibility.

    Args:
        output_dir: Output directory containing 0-deck-sections/
        screenshots: List of screenshot metadata dicts with path, category, description

    Returns:
        Number of sections updated with screenshots
    """
    if not screenshots:
        return 0

    deck_sections_dir = output_dir / "0-deck-sections"
    if not deck_sections_dir.exists():
        print(f"  ⚠️  0-deck-sections/ not found, cannot embed screenshots", flush=True)
        return 0

    # Group screenshots by target section
    section_to_screenshots: Dict[str, List[Dict]] = {}
    for screenshot in screenshots:
        category = screenshot.get("category", "general")
        target_section = SCREENSHOT_CATEGORY_TO_SECTION.get(category)

        if target_section:
            if target_section not in section_to_screenshots:
                section_to_screenshots[target_section] = []
            section_to_screenshots[target_section].append(screenshot)

    if not section_to_screenshots:
        print(f"  No screenshots mapped to sections", flush=True)
        return 0

    updated_count = 0

    for section_filename, section_screenshots in section_to_screenshots.items():
        section_path = deck_sections_dir / section_filename

        if not section_path.exists():
            print(f"  ⚠️  Section file {section_filename} not found, skipping screenshot embed", flush=True)
            continue

        # Read existing content
        existing_content = section_path.read_text()

        # Build image embeds (use absolute paths for cross-directory compatibility)
        image_embeds = []
        for ss in section_screenshots:
            # Build absolute path to screenshot
            abs_screenshot_path = output_dir / ss["path"]
            description = ss.get("description", f"Page {ss.get('page_number', '?')} from pitch deck")

            # Create markdown image embed
            image_embeds.append(f"![{description}]({abs_screenshot_path})")

        # Prepend images to content
        images_block = "\n\n".join(image_embeds)
        new_content = f"{images_block}\n\n{existing_content}"

        # Write updated content
        section_path.write_text(new_content)
        updated_count += 1
        print(f"  ✓ Embedded {len(section_screenshots)} screenshot(s) in {section_filename}", flush=True)

    return updated_count


def deck_analyst_agent(state: Dict) -> Dict:
    """
    Analyzes pitch deck and extracts key information.

    CRITICAL: Only handles PDF decks for now. Image decks cause bottlenecks.
    Future: Add image deck support with optimization (resizing, compression).

    Args:
        state: Current memo state

    Returns:
        Updated state with deck_analysis and initial section drafts
    """
    deck_path = state.get("deck_path")

    if not deck_path or not Path(deck_path).exists():
        return {
            "deck_analysis": None,
            "messages": ["No deck available, skipping deck analysis"]
        }

    deck_file = Path(deck_path)

    # STEP 1: Extract text from deck (PDF or PowerPoint)
    deck_suffix = deck_file.suffix.lower()

    if deck_suffix == ".pdf":
        print(f"Extracting text from PDF deck ({deck_file.name})...", flush=True)
        deck_content = extract_text_from_pdf(deck_path)
        print(f"Extracted {len(deck_content)} characters from deck", flush=True)

        # If minimal text extracted, it's an image-based PDF - use Claude's PDF vision
        # Threshold of 1000 chars accounts for page headers (22 pages × ~17 chars = ~374 chars)
        # A text-based deck would have substantially more content
        if len(deck_content.strip()) < 1000:
            print("⚠️  Minimal text extracted - PDF appears to be image-based", flush=True)
            print("Using Claude's PDF vision to analyze deck...", flush=True)
            return analyze_pdf_with_vision(deck_path, state)

    elif deck_suffix in [".pptx", ".ppt"]:
        print(f"Extracting text from PowerPoint deck ({deck_file.name})...", flush=True)
        try:
            deck_content = extract_text_from_pptx(deck_path)
            print(f"Extracted {len(deck_content)} characters from {len(deck_content.split('--- SLIDE'))-1} slides", flush=True)

            if len(deck_content.strip()) < 500:
                print("⚠️  Minimal text extracted - PowerPoint may be mostly images", flush=True)
                print("Note: Image-only PowerPoint analysis not yet supported", flush=True)
                # Continue with what we have
        except Exception as e:
            return {
                "deck_analysis": None,
                "messages": [f"Error reading PowerPoint: {e}"]
            }
    else:
        # Unsupported format
        return {
            "deck_analysis": None,
            "messages": [f"Deck format {deck_suffix} not supported. Supported: .pdf, .pptx"]
        }

    # STEP 2: Analyze extracted text with Claude
    print("Analyzing deck content with Claude Sonnet 4.5...", flush=True)
    llm = ChatAnthropic(
        model="claude-sonnet-4-5-20250929",
        temperature=0,
        timeout=120  # 2 minute timeout to avoid hangs
    )

    analysis_prompt = f"""You are a venture capital investment analyst reviewing a pitch deck.

PITCH DECK CONTENT:
{deck_content}

Extract the following information in JSON format:
{{
  "company_name": "...",
  "tagline": "...",
  "problem_statement": "...",
  "solution_description": "...",
  "product_description": "...",
  "business_model": "...",
  "market_size": {{"TAM": "...", "SAM": "...", "SOM": "..."}},
  "traction_metrics": [{{"metric": "...", "value": "..."}}, ...],
  "team_members": [{{"name": "...", "role": "...", "background": "..."}}, ...],
  "funding_ask": "...",
  "use_of_funds": ["...", "..."],
  "competitive_landscape": "...",
  "go_to_market": "...",
  "milestones": ["...", "..."],
  "extraction_notes": ["List what info was found vs. missing"]
}}

IMPORTANT:
- Only include information explicitly stated in the deck
- Use "Not mentioned" if information is absent
- Capture specific numbers (revenue, users, growth rates)
- Note the deck's strengths and weaknesses in extraction_notes
- Return ONLY valid JSON, no other text
"""

    print("Sending deck to Claude for analysis...", flush=True)
    response = llm.invoke(analysis_prompt)
    print("Deck analysis complete, parsing results...", flush=True)

    # Parse JSON from response
    try:
        deck_analysis = json.loads(response.content)
    except json.JSONDecodeError:
        # Try to extract JSON from markdown code block
        print("Extracting JSON from markdown code block...", flush=True)
        content = response.content
        if "```json" in content:
            json_start = content.find("```json") + 7
            json_end = content.find("```", json_start)
            content = content[json_start:json_end].strip()
        elif "```" in content:
            json_start = content.find("```") + 3
            json_end = content.find("```", json_start)
            content = content[json_start:json_end].strip()
        deck_analysis = json.loads(content)

    # Get page/slide count based on file type
    if deck_suffix == ".pdf":
        deck_analysis["deck_page_count"] = len(PdfReader(deck_path).pages)
    elif deck_suffix in [".pptx", ".ppt"]:
        deck_analysis["deck_page_count"] = len(Presentation(deck_path).slides)
    else:
        deck_analysis["deck_page_count"] = 0

    print(f"Extracted data from {deck_analysis['deck_page_count']}-page deck", flush=True)

    # STEP 3: Create initial section drafts where relevant info exists
    print("Creating initial section drafts from deck data...", flush=True)
    section_drafts = create_initial_section_drafts(deck_analysis, state, llm)
    print(f"Created {len(section_drafts)} initial section drafts", flush=True)

    # STEP 4: Save artifacts
    from ..artifacts import save_deck_analysis_artifacts
    print("Saving deck analysis artifacts...", flush=True)
    # Extract just the content strings for artifact saving
    section_drafts_for_disk = {k: v["content"] for k, v in section_drafts.items()}
    # Pass firm from state for firm-scoped output paths
    firm = state.get("firm")
    # Use output_dir from state (created at workflow start)
    state_output_dir = state.get("output_dir")
    output_dir = save_deck_analysis_artifacts(
        state["company_name"],
        deck_analysis,
        section_drafts_for_disk,
        firm=firm,
        output_dir=Path(state_output_dir) if state_output_dir else None
    )

    # STEP 5: Extract visual screenshots (async with timeout, non-blocking on failure)
    deck_screenshots = []
    if deck_suffix == ".pdf":
        deck_screenshots = extract_screenshots_with_timeout(
            deck_path,
            deck_analysis,
            output_dir,  # Pass the exact output_dir from artifact saving
            timeout_seconds=90,  # 90 second timeout for screenshot extraction
            firm=firm
        )
        if deck_screenshots:
            deck_analysis["screenshots"] = deck_screenshots

            # STEP 6: Embed screenshots into section files
            # This updates 0-deck-sections/ files with image embeds at the top
            print("Embedding screenshots in section drafts...", flush=True)
            sections_updated = embed_screenshots_in_section_files(output_dir, deck_screenshots)
            if sections_updated:
                print(f"✓ Embedded screenshots in {sections_updated} section file(s)", flush=True)

            # STEP 7: Build and save section-to-screenshots mapping
            # This mapping tells inject_deck_images agent which images go in which sections
            print("Building section-to-screenshots mapping...", flush=True)
            screenshot_mapping = build_section_to_screenshots(deck_screenshots, output_dir)
            deck_analysis["section_to_screenshots"] = screenshot_mapping
            print(f"✓ Mapped {len(screenshot_mapping.get('images', []))} screenshots for injection", flush=True)

            # Re-save deck_analysis.json with screenshots and mappings
            deck_analysis_path = output_dir / "0-deck-analysis.json"
            with open(deck_analysis_path, "w") as f:
                json.dump(deck_analysis, f, indent=2, default=str)
            print(f"✓ Updated 0-deck-analysis.json with screenshot mappings", flush=True)

    return {
        "deck_analysis": deck_analysis,
        "draft_sections": section_drafts,  # Partial sections
        "messages": [f"Deck analysis complete: {deck_analysis['deck_page_count']} pages analyzed, {len(deck_screenshots)} screenshots extracted"]
    }


def extract_text_from_pdf(pdf_path: str) -> str:
    """Extract text from PDF using pypdf."""
    reader = PdfReader(pdf_path)
    text_content = []

    for page_num, page in enumerate(reader.pages, 1):
        page_text = page.extract_text()
        text_content.append(f"--- PAGE {page_num} ---\n{page_text}\n")

    return "\n".join(text_content)


def extract_text_from_pptx(pptx_path: str) -> str:
    """
    Extract text from PowerPoint file.

    Extracts text from all slides, including:
    - Slide titles
    - Text boxes and shapes
    - Tables
    - Notes

    Args:
        pptx_path: Path to .pptx file

    Returns:
        Extracted text with slide markers
    """
    prs = Presentation(pptx_path)
    text_content = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        slide_text.append(f"--- SLIDE {slide_num} ---")

        # Extract from all shapes
        for shape in slide.shapes:
            # Handle text frames (most common)
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = "".join([run.text for run in paragraph.runs])
                    if para_text.strip():
                        slide_text.append(para_text.strip())

            # Handle tables
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        slide_text.append(" | ".join(row_text))

        # Extract slide notes if present
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_text.append(f"[Notes: {notes}]")

        text_content.append("\n".join(slide_text))

    return "\n\n".join(text_content)


def extract_deck_screenshots(
    pdf_path: str,
    output_dir: Path,
    page_selections: List[Dict[str, Any]] = None,
    use_pdf2image: bool = True,
    dpi: int = 150,
    quality: int = 85,
    max_width: int = 1200
) -> List[Dict[str, Any]]:
    """
    Extract screenshots from specific PDF pages.

    Args:
        pdf_path: Path to PDF file
        output_dir: Directory to save screenshots
        page_selections: List of dicts with page_number, category, and description
                        If None, extracts all pages
        use_pdf2image: Use pdf2image (higher quality) vs PyMuPDF (faster)
        dpi: Resolution for rendering (150 = good balance of quality/size)
        quality: JPEG quality (1-100)
        max_width: Maximum width in pixels (default 1200). Images wider than this
                   are resized while maintaining aspect ratio.

    Returns:
        List of dicts with path, page_number, category, description, dimensions
    """
    screenshots_dir = output_dir / "deck-screenshots"
    screenshots_dir.mkdir(exist_ok=True)

    extracted = []
    doc = fitz.open(pdf_path)
    page_count = len(doc)

    # If no selections provided, we won't extract anything by default
    # (LLM should guide which pages to extract)
    if not page_selections:
        doc.close()
        return extracted

    # Extract selected pages
    for selection in page_selections:
        page_num = selection.get("page_number", 1) - 1  # Convert to 0-indexed
        if page_num < 0 or page_num >= page_count:
            continue

        category = selection.get("category", "general")
        slug = selection.get("slug", "")
        description = selection.get("description", f"Page {page_num + 1}")

        # Generate filename: page-03-team-founding-team.png
        safe_category = "".join(c for c in category if c.isalnum() or c == '-').lower()
        safe_slug = "".join(c for c in slug if c.isalnum() or c == '-').lower()
        if safe_slug:
            filename = f"page-{page_num + 1:02d}-{safe_category}-{safe_slug}.png"
        else:
            filename = f"page-{page_num + 1:02d}-{safe_category}.png"
        output_path = screenshots_dir / filename

        try:
            img = None
            if use_pdf2image and PDF2IMAGE_AVAILABLE:
                # Higher quality rendering via Poppler
                images = convert_from_path(
                    pdf_path,
                    dpi=dpi,
                    first_page=page_num + 1,
                    last_page=page_num + 1,
                    fmt="png"
                )
                if images:
                    img = images[0]

            # Fallback to PyMuPDF rendering if pdf2image failed or unavailable
            if img is None:
                page = doc[page_num]
                # Scale factor: 150 DPI / 72 base DPI = ~2.08x
                scale = dpi / 72.0
                mat = fitz.Matrix(scale, scale)
                pix = page.get_pixmap(matrix=mat)

                # Convert to PIL for saving as PNG
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

            # Resize if width exceeds max_width (maintain aspect ratio)
            width, height = img.size
            if width > max_width:
                ratio = max_width / width
                new_height = int(height * ratio)
                img = img.resize((max_width, new_height), Image.Resampling.LANCZOS)
                width, height = img.size
                print(f"    📐 Resized to {width}x{height} (max width: {max_width}px)")

            # Save with optimization
            img.save(output_path, "PNG", optimize=True)

            extracted.append({
                "path": str(output_path.relative_to(output_dir)),
                "filename": filename,
                "page_number": page_num + 1,
                "category": category,
                "slug": slug,
                "description": description,
                "width": width,
                "height": height
            })
            print(f"    ✓ Extracted page {page_num + 1}: {filename}")

        except Exception as e:
            print(f"    ⚠ Failed to extract page {page_num + 1}: {e}")

    doc.close()
    return extracted


def identify_visual_pages(
    pdf_path: str,
    deck_analysis: Dict[str, Any],
    client: Anthropic,
    firm: str = None
) -> List[Dict[str, Any]]:
    """
    Use Claude to identify which pages contain valuable visual content.

    Loads the deck classification guide from templates/deck-classification-guide.md
    and passes it to Claude Vision for accurate, granular slide classification.

    Args:
        pdf_path: Path to PDF file
        deck_analysis: Already extracted deck analysis data
        client: Anthropic client for vision API
        firm: Optional firm name for firm-scoped classification guide

    Returns:
        List of page selections with category, slug, and description
    """
    doc = fitz.open(pdf_path)
    page_count = len(doc)

    # Sample pages for analysis (all pages if <= 20, otherwise sample)
    if page_count <= 20:
        pages_to_analyze = list(range(page_count))
    else:
        # Sample: first 5, last 5, and evenly distributed middle pages
        pages_to_analyze = list(range(5)) + list(range(page_count - 5, page_count))
        middle_count = min(10, page_count - 10)
        step = (page_count - 10) // middle_count
        pages_to_analyze += [5 + i * step for i in range(middle_count)]
        pages_to_analyze = sorted(set(pages_to_analyze))

    # Convert pages to images for Claude
    image_contents = []
    for page_num in pages_to_analyze:
        page = doc[page_num]
        # 0.5x resolution — readable enough to classify tables, charts, and text
        mat = fitz.Matrix(0.5, 0.5)
        pix = page.get_pixmap(matrix=mat)

        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        img_buffer = BytesIO()
        img.save(img_buffer, format="JPEG", quality=70)
        img_bytes = img_buffer.getvalue()
        img_b64 = base64.standard_b64encode(img_bytes).decode("utf-8")

        image_contents.append({
            "type": "image",
            "source": {
                "type": "base64",
                "media_type": "image/jpeg",
                "data": img_b64
            }
        })

    doc.close()

    # Load classification guide
    classification_guide = load_classification_guide(firm=firm)

    if classification_guide:
        # Use the full classification guide as context
        prompt = f"""Analyze these {len(pages_to_analyze)} pitch deck slides and identify which ones contain valuable VISUAL content that should be saved as screenshots for an investment memo.

The slides are numbered in order: {[p + 1 for p in pages_to_analyze]}

Use the following classification guide to categorize each slide. Read it carefully — it defines the exact categories, what to look for, what NOT to classify as each category, and how to resolve ambiguity.

---
{classification_guide}
---

Return a JSON array of pages to extract. Only include pages with SIGNIFICANT visual value — skip title slides, text-only slides, agendas, and legal disclaimers (see "Slides to Skip" in the guide above).

Each entry MUST include "category" (from the guide), "slug" (2-5 word descriptor, lowercase hyphenated, unique across the deck), and "description" (one sentence about the specific visual content).

Format:
```json
[
  {{"page_number": 3, "category": "team", "slug": "founding-team", "description": "Founding team photos with bios and university/company backgrounds"}},
  {{"page_number": 7, "category": "traction", "slug": "revenue-growth", "description": "MRR growth chart showing 3x YoY growth from $50K to $150K"}},
  {{"page_number": 14, "category": "unit-economics", "slug": "unit-margins", "description": "Per-unit COGS and margin table showing 72-87% gross margins"}}
]
```

Return ONLY the JSON array, no other text. If no pages have significant visual value, return an empty array: []"""
    else:
        # Fallback: inline categories if guide file is missing
        prompt = f"""Analyze these {len(pages_to_analyze)} pitch deck slides and identify which ones contain valuable VISUAL content that should be saved as screenshots for an investment memo.

The slides are numbered in order: {[p + 1 for p in pages_to_analyze]}

Classify each slide into one of these categories:
- **overview**: Company overview, mission, elevator pitch
- **problem**: Market/industry problems, status quo failures
- **customer-pain**: Specific user frustrations, persona pain points
- **ideal-customer-profile**: Target personas, customer segments, buyer demographics
- **solution**: How the company solves the problem (conceptual)
- **product-demo**: Product UI screenshots, app interfaces, physical product photos
- **value-proposition**: Key benefits, feature highlights, "why us"
- **technology**: Architecture diagrams, science mechanisms, IP/patents, R&D
- **business-model**: Revenue model, pricing, monetization strategy
- **unit-economics**: Margins, COGS, LTV/CAC, per-unit profitability
- **market-size**: TAM/SAM/SOM, market growth charts
- **competitive-positioning**: 2x2 matrices, feature comparison grids
- **competition-landscape**: Competitor lists, market maps
- **traction**: Growth metrics, revenue charts, milestones
- **team**: Team photos, org charts, advisor headshots
- **go-to-market**: GTM strategy, acquisition channels, funnels
- **fundraising**: Round details, use of funds, investor logos
- **financials**: Revenue projections, P&L forecasts, runway
- **partnerships**: Strategic partners, ecosystem diagrams
- **branding**: Brand positioning, packaging, D2C website
- **vision**: Big-picture roadmap, long-term vision
- **impact**: Social impact, ESG, sustainability

Return a JSON array of pages to extract. Only include pages with SIGNIFICANT visual value (not just text slides or bullet points).

Each entry MUST include "category", "slug" (2-5 word descriptor, lowercase hyphenated), and "description".

Format:
```json
[
  {{"page_number": 3, "category": "team", "slug": "founding-team", "description": "Founding team photos with backgrounds"}},
  {{"page_number": 7, "category": "traction", "slug": "revenue-growth", "description": "MRR growth chart showing 3x YoY growth"}},
  {{"page_number": 14, "category": "unit-economics", "slug": "unit-margins", "description": "Per-unit margin table with COGS breakdown"}}
]
```

Return ONLY the JSON array, no other text. If no pages have significant visual value, return an empty array: []"""

    try:
        content_blocks = image_contents + [{"type": "text", "text": prompt}]

        response = client.messages.create(
            model="claude-sonnet-4-5-20250929",
            max_tokens=4000,
            messages=[{
                "role": "user",
                "content": content_blocks
            }]
        )

        content = response.content[0].text

        # Parse JSON response
        try:
            selections = json.loads(content)
        except json.JSONDecodeError:
            # Extract from code block if needed
            if "```json" in content:
                json_str = content.split("```json")[1].split("```")[0].strip()
            elif "```" in content:
                json_str = content.split("```")[1].split("```")[0].strip()
            else:
                json_str = content
            selections = json.loads(json_str)

        # Validate and clean selections
        valid_selections = []
        for sel in selections:
            if isinstance(sel, dict) and "page_number" in sel:
                valid_selections.append({
                    "page_number": int(sel.get("page_number", 1)),
                    "category": str(sel.get("category", "general")),
                    "slug": str(sel.get("slug", ""))[:50],
                    "description": str(sel.get("description", ""))[:200]
                })

        return valid_selections

    except Exception as e:
        print(f"    ⚠ Visual page identification failed: {e}")
        return []


def analyze_pdf_with_vision(pdf_path: str, state: Dict) -> Dict:
    """
    Analyze image-based PDF by extracting images and sending to Claude's vision API.

    Args:
        pdf_path: Path to PDF file
        state: Current memo state

    Returns:
        Updated state with deck analysis
    """
    # Extract images from PDF using PyMuPDF
    print(f"Extracting images from PDF...", flush=True)
    doc = fitz.open(pdf_path)
    page_count = len(doc)

    # Process in batches of 5 to avoid API payload limits
    batch_size = 5
    all_deck_analyses = []

    # Use Anthropic client directly for vision
    client = Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))

    for batch_start in range(0, page_count, batch_size):
        batch_end = min(batch_start + batch_size, page_count)
        batch_num = (batch_start // batch_size) + 1
        total_batches = (page_count + batch_size - 1) // batch_size

        print(f"Processing slides {batch_start + 1}-{batch_end} (batch {batch_num}/{total_batches})...", flush=True)

        # Convert batch pages to images
        image_contents = []
        for page_num in range(batch_start, batch_end):
            page = doc[page_num]
            # Render at 0.5x scale to reduce payload size (still readable)
            mat = fitz.Matrix(0.5, 0.5)
            pix = page.get_pixmap(matrix=mat)

            # Convert to PIL Image to compress as JPEG
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            img_buffer = BytesIO()
            img.save(img_buffer, format="JPEG", quality=85, optimize=True)
            img_bytes = img_buffer.getvalue()
            img_b64 = base64.standard_b64encode(img_bytes).decode("utf-8")

            image_contents.append({
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": "image/jpeg",
                    "data": img_b64
                }
            })

        print(f"Sending batch {batch_num} ({len(image_contents)} slides) to Claude...", flush=True)

        analysis_prompt = f"""You are a venture capital investment analyst reviewing slides {batch_start + 1}-{batch_end} of a pitch deck.

Extract the following information from these slides in JSON format (use "Not mentioned" for information not present in THIS batch):
{{
  "company_name": "...",
  "tagline": "...",
  "problem_statement": "...",
  "solution_description": "...",
  "product_description": "...",
  "business_model": "...",
  "market_size": {{"TAM": "...", "SAM": "...", "SOM": "..."}},
  "traction_metrics": [{{"metric": "...", "value": "..."}}, ...],
  "team_members": [{{"name": "...", "role": "...", "background": "..."}}, ...],
  "funding_ask": "...",
  "use_of_funds": ["...", "..."],
  "competitive_landscape": "...",
  "go_to_market": "...",
  "milestones": ["...", "..."],
  "extraction_notes": ["List what info was found vs. missing"]
}}

IMPORTANT:
- Only include information explicitly stated in the deck
- Use "Not mentioned" if information is absent
- Capture specific numbers (revenue, users, growth rates, funding amounts)
- Note the deck's strengths and weaknesses in extraction_notes
- Return ONLY valid JSON, no other text"""

        try:
            # Build content blocks: all images followed by the prompt
            content_blocks = image_contents + [{"type": "text", "text": analysis_prompt}]

            response = client.messages.create(
                model="claude-sonnet-4-5-20250929",
                max_tokens=4000,
                messages=[{
                    "role": "user",
                    "content": content_blocks
                }]
            )

            print(f"✓ Batch {batch_num} complete, parsing results...", flush=True)

            # Extract JSON from response
            content = response.content[0].text

            try:
                batch_analysis = json.loads(content)
                print(f"✓ Batch {batch_num}: Parsed JSON directly", flush=True)
            except json.JSONDecodeError as e:
                # Try to extract JSON from markdown code block
                print(f"Batch {batch_num}: Direct JSON parse failed, extracting from code block...", flush=True)
                if "```json" in content:
                    json_str = content.split("```json")[1].split("```")[0].strip()
                elif "```" in content:
                    json_str = content.split("```")[1].split("```")[0].strip()
                else:
                    json_str = content

                try:
                    batch_analysis = json.loads(json_str)
                    print(f"✓ Batch {batch_num}: Parsed JSON from code block", flush=True)
                except json.JSONDecodeError as e2:
                    print(f"ERROR: Batch {batch_num} failed to parse JSON: {e2}", flush=True)
                    print(f"First 500 chars of response: {content[:500]}", flush=True)
                    # Continue with next batch instead of failing completely
                    continue

            all_deck_analyses.append(batch_analysis)
            print(f"✓ Batch {batch_num}: Found {len(batch_analysis.get('team_members', []))} team members, "
                  f"{len(batch_analysis.get('traction_metrics', []))} traction metrics", flush=True)

        except Exception as e:
            print(f"ERROR: Batch {batch_num} failed: {e}", flush=True)
            # Continue with next batch instead of failing completely
            continue

    # Close the PDF document
    doc.close()

    # Merge all batch analyses into a single comprehensive analysis
    print(f"\nMerging {len(all_deck_analyses)} batch analyses...", flush=True)

    if not all_deck_analyses:
        print("ERROR: No batches were successfully analyzed", flush=True)
        return {
            "deck_analysis": None,
            "messages": ["Deck analysis failed (vision mode): No batches succeeded"]
        }

    # Start with the first batch as base
    deck_analysis = all_deck_analyses[0].copy()
    deck_analysis["deck_page_count"] = page_count

    # Merge data from subsequent batches
    for batch_data in all_deck_analyses[1:]:
        # For string fields: use first non-"Not mentioned" value
        for field in ["company_name", "tagline", "problem_statement", "solution_description",
                      "product_description", "business_model", "funding_ask", "competitive_landscape", "go_to_market"]:
            if batch_data.get(field) and batch_data[field] != "Not mentioned":
                if deck_analysis.get(field) == "Not mentioned" or not deck_analysis.get(field):
                    deck_analysis[field] = batch_data[field]

        # For dict fields (market_size): merge non-"Not mentioned" values
        if batch_data.get("market_size"):
            if not deck_analysis.get("market_size"):
                deck_analysis["market_size"] = {}
            for key, value in batch_data["market_size"].items():
                if value and value != "Not mentioned":
                    if not deck_analysis["market_size"].get(key) or deck_analysis["market_size"][key] == "Not mentioned":
                        deck_analysis["market_size"][key] = value

        # For list fields: append unique items
        for field in ["traction_metrics", "team_members", "use_of_funds", "milestones", "extraction_notes"]:
            if batch_data.get(field) and isinstance(batch_data[field], list):
                if not deck_analysis.get(field):
                    deck_analysis[field] = []
                # Append items that aren't already present
                for item in batch_data[field]:
                    if item not in deck_analysis[field] and item != "Not mentioned":
                        deck_analysis[field].append(item)

    print(f"✓ Merged analysis complete: {deck_analysis.get('deck_page_count', 'unknown')} pages analyzed", flush=True)
    print(f"✓ Final totals: {len(deck_analysis.get('team_members', []))} team members, "
          f"{len(deck_analysis.get('traction_metrics', []))} traction metrics", flush=True)

    # Continue with the same flow as text-based analysis
    try:
        llm = ChatAnthropic(
            model="claude-sonnet-4-5-20250929",
            temperature=0
        )

        print("Creating initial section drafts from deck data...", flush=True)
        section_drafts = create_initial_section_drafts(deck_analysis, state, llm)
        print(f"Created {len(section_drafts)} initial section drafts", flush=True)

        # Save artifacts (same as text-based path)
        from ..artifacts import save_deck_analysis_artifacts
        print("Saving deck analysis artifacts...", flush=True)
        # Extract just the content strings for artifact saving
        section_drafts_for_disk = {k: v["content"] for k, v in section_drafts.items()}
        # Pass firm from state for firm-scoped output paths
        firm = state.get("firm")
        # Use output_dir from state (created at workflow start)
        state_output_dir = state.get("output_dir")
        save_deck_analysis_artifacts(
            state["company_name"],
            deck_analysis,
            section_drafts_for_disk,
            firm=firm,
            output_dir=Path(state_output_dir) if state_output_dir else None
        )

        # Extract visual screenshots (vision mode already has client)
        deck_screenshots = []
        print("Identifying visual pages for screenshot extraction...", flush=True)
        try:
            from ..utils import get_output_dir_from_state

            output_dir = get_output_dir_from_state(state)

            # Reuse the client we already have
            page_selections = identify_visual_pages(pdf_path, deck_analysis, client, firm=firm)

            if page_selections:
                print(f"  Found {len(page_selections)} pages with visual content", flush=True)
                deck_screenshots = extract_deck_screenshots(
                    pdf_path,
                    output_dir,
                    page_selections,
                    use_pdf2image=PDF2IMAGE_AVAILABLE,
                    dpi=150
                )
                print(f"  Extracted {len(deck_screenshots)} screenshots", flush=True)
                deck_analysis["screenshots"] = deck_screenshots

                # Embed screenshots into section files
                print("Embedding screenshots in section drafts...", flush=True)
                sections_updated = embed_screenshots_in_section_files(output_dir, deck_screenshots)
                if sections_updated:
                    print(f"✓ Embedded screenshots in {sections_updated} section file(s)", flush=True)

                # Build and save section-to-screenshots mapping
                print("Building section-to-screenshots mapping...", flush=True)
                screenshot_mapping = build_section_to_screenshots(deck_screenshots, output_dir)
                deck_analysis["section_to_screenshots"] = screenshot_mapping
                print(f"✓ Mapped {len(screenshot_mapping.get('images', []))} screenshots for injection", flush=True)

                # Re-save deck_analysis.json with screenshots and mappings
                deck_analysis_path = output_dir / "0-deck-analysis.json"
                with open(deck_analysis_path, "w") as f:
                    json.dump(deck_analysis, f, indent=2, default=str)
                print(f"✓ Updated 0-deck-analysis.json with screenshot mappings", flush=True)
            else:
                print("  No significant visual content identified", flush=True)

        except Exception as e:
            print(f"  ⚠ Screenshot extraction failed: {e}", flush=True)

        return {
            "deck_analysis": deck_analysis,
            "draft_sections": section_drafts,
            "messages": [f"Deck analysis complete (vision mode): {deck_analysis.get('deck_page_count', 'unknown')} pages analyzed, {len(deck_screenshots)} screenshots extracted"]
        }

    except Exception as e:
        print(f"Error analyzing PDF with vision: {e}", flush=True)
        return {
            "deck_analysis": None,
            "messages": [f"Deck analysis failed (vision mode): {str(e)}"]
        }


def create_initial_section_drafts(deck_analysis: Dict, state: Dict, llm: ChatAnthropic) -> Dict[str, SectionDraft]:
    """
    Create draft sections for ALL extracted deck data fields.

    This creates a draft for EVERY field that has substantial data, regardless of
    outline section names. Downstream agents (researcher, writer) can then use
    these drafts as authoritative source material.

    Args:
        deck_analysis: Extracted deck data
        state: Current memo state
        llm: Language model for generating drafts

    Returns:
        Dictionary mapping section filenames to SectionDraft objects
    """
    drafts = {}

    # Map deck fields to filenames - outline agnostic
    # Create a draft for EVERY field with data, let downstream agents use them
    DECK_FIELD_CONFIG = {
        "problem_statement": {
            "filename": "deck-problem.md",
            "display_name": "Problem Statement",
            "related_fields": ["problem_statement"]
        },
        "solution_description": {
            "filename": "deck-solution.md",
            "display_name": "Solution",
            "related_fields": ["solution_description"]
        },
        "product_description": {
            "filename": "deck-product.md",
            "display_name": "Product",
            "related_fields": ["product_description"]
        },
        "business_model": {
            "filename": "deck-business-model.md",
            "display_name": "Business Model",
            "related_fields": ["business_model"]
        },
        "market_size": {
            "filename": "deck-market.md",
            "display_name": "Market Size",
            "related_fields": ["market_size"]
        },
        "competitive_landscape": {
            "filename": "deck-competitive.md",
            "display_name": "Competitive Landscape",
            "related_fields": ["competitive_landscape"]
        },
        "traction_metrics": {
            "filename": "deck-traction.md",
            "display_name": "Traction & Metrics",
            "related_fields": ["traction_metrics", "milestones"]
        },
        "team_members": {
            "filename": "deck-team.md",
            "display_name": "Team",
            "related_fields": ["team_members"]
        },
        "funding_ask": {
            "filename": "deck-funding.md",
            "display_name": "Funding & Terms",
            "related_fields": ["funding_ask", "use_of_funds"]
        },
        "go_to_market": {
            "filename": "deck-gtm.md",
            "display_name": "Go-to-Market Strategy",
            "related_fields": ["go_to_market"]
        },
    }

    def has_substantial_data(field_value) -> bool:
        """Check if field has real data (not empty or placeholder)."""
        if not field_value:
            return False
        if field_value == "Not mentioned":
            return False
        if field_value == "":
            return False
        if field_value == []:
            return False
        if field_value == {}:
            return False
        # For dicts, check if all values are "Not mentioned"
        if isinstance(field_value, dict):
            return any(v and v != "Not mentioned" for v in field_value.values())
        return True

    # Create a draft for each field with substantial data
    for primary_field, config in DECK_FIELD_CONFIG.items():
        # Check if any of the related fields have data
        related_fields = config["related_fields"]
        fields_with_data = [
            f for f in related_fields
            if has_substantial_data(deck_analysis.get(f))
        ]

        if not fields_with_data:
            continue

        # Generate draft content
        content = create_section_draft_from_deck(
            llm,
            config["display_name"],
            deck_analysis,
            fields_with_data
        )

        # Create SectionDraft object
        drafts[config["filename"]] = SectionDraft(
            section_name=config["display_name"],
            content=content,
            word_count=len(content.split()),
            citations=[]  # Citations added by citation enrichment agent
        )

        print(f"    ✓ Created deck draft: {config['filename']} ({len(content.split())} words)", flush=True)

    return drafts


def create_section_draft_from_deck(llm: ChatAnthropic, section_name: str, deck_data: Dict, fields: List[str]) -> str:
    """
    Generate a section draft from deck data.

    Args:
        llm: Language model for generating drafts
        section_name: Name of the section being drafted
        deck_data: Extracted deck data
        fields: Relevant fields for this section

    Returns:
        Draft section content in markdown
    """
    relevant_data = {k: deck_data.get(k) for k in fields if deck_data.get(k)}

    prompt = f"""Draft the "{section_name}" section for an investment memo based on this pitch deck data:

{json.dumps(relevant_data, indent=2)}

Write a concise, analytical section (200-400 words) that:
- Uses specific numbers and metrics from the deck
- Maintains analytical (not promotional) tone
- Notes data gaps explicitly (e.g., "Team backgrounds not disclosed in deck")
- Formats for readability (bullet points where appropriate)

This is an INITIAL DRAFT. The Research and Writer agents will augment with external data.

Return ONLY the section content in markdown format, no preamble.
"""

    response = llm.invoke(prompt)
    return response.content
